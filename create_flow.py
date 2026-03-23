from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(30, 39, 97)       # Navy blue
ACCENT = RGBColor(245, 158, 11)     # Orange
LIGHT = RGBColor(248, 250, 252)     # Off-white
TEXT = RGBColor(30, 41, 59)         # Dark slate

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = ACCENT
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY

    return slide

def add_box(slide, x, y, w, h, text, bg_color, border_color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_left = Inches(0.15)
    text_frame.margin_right = Inches(0.15)

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 255, 255) if bg_color != LIGHT else TEXT
    p.alignment = PP_ALIGN.LEFT  # Left

# Slide 1: Title
add_title_slide(prs, "PhotoFinder", "Face Recognition Photo Matching System")

# Slide 2: Architecture
slide = add_content_slide(prs, "System Architecture")

add_box(slide, 0.5, 1.2, 2.8, 1.8, "📱 Admin Dashboard\n\n• Submit Drive\n• Create Event\n• Download QR", PRIMARY)
add_box(slide, 3.7, 1.2, 2.8, 1.8, "⚙️ Backend Engine\n\n• DeepFace AI\n• Auto-Sync\n• Face Encoding", ACCENT)
add_box(slide, 6.9, 1.2, 2.8, 1.8, "👤 User Portal\n\n• Scan QR\n• Upload Selfie\n• View Matches", PRIMARY)

add_box(slide, 3.1, 3.5, 3.8, 0.9, "💾 Storage: Google Drive + Local Encodings", RGBColor(226, 232, 240))

# Slide 3: Admin Workflow
slide = add_content_slide(prs, "Admin Workflow")

steps = [
    (0.5, 1.2, "1", "Submit Drive Link"),
    (2.5, 1.2, "2", "Download Photos"),
    (4.5, 1.2, "3", "Encode Faces"),
    (6.5, 1.2, "4", "Generate QR"),
    (0.5, 2.8, "5", "Status Ready"),
    (2.5, 2.8, "6", "Share QR Code")
]

for x, y, num, title in steps:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.color.rgb = PRIMARY

    text_frame = circle.text_frame
    p = text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = 1

    label_box = slide.shapes.add_textbox(Inches(x + 0.7), Inches(y + 0.05), Inches(1.5), Inches(0.4))
    label_frame = label_box.text_frame
    label_frame.text = title
    label_frame.paragraphs[0].font.size = Pt(10)
    label_frame.paragraphs[0].font.bold = True
    label_frame.paragraphs[0].font.color.rgb = PRIMARY

# Slide 4: User Workflow
slide = add_content_slide(prs, "User Workflow")

steps = [
    (0.5, 1.2, "1", "Scan QR"),
    (2.5, 1.2, "2", "Open Event"),
    (4.5, 1.2, "3", "Upload Selfie"),
    (6.5, 1.2, "4", "Face Match"),
    (0.5, 2.8, "5", "View Results"),
    (2.5, 2.8, "6", "Download ZIP")
]

for x, y, num, title in steps:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.color.rgb = PRIMARY

    text_frame = circle.text_frame
    p = text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = 1

    label_box = slide.shapes.add_textbox(Inches(x + 0.7), Inches(y + 0.05), Inches(1.5), Inches(0.4))
    label_frame = label_box.text_frame
    label_frame.text = title
    label_frame.paragraphs[0].font.size = Pt(10)
    label_frame.paragraphs[0].font.bold = True
    label_frame.paragraphs[0].font.color.rgb = PRIMARY

# Slide 5: Auto-Sync & Real-Time Updates
slide = add_content_slide(prs, "Real-Time Auto-Sync & Updates ⚡")

# Admin auto-refresh
admin_text = "🔄 Admin Dashboard\n\nAuto-Refresh Every 10 Minutes:\n\n✓ Check Drive for new photos\n✓ Auto-load in background\n✓ Update photo count\n✓ Re-index automatically"
add_box(slide, 0.5, 1.1, 4.2, 2.8, admin_text, RGBColor(224, 231, 255))

# User auto-refresh
user_text = "📱 User Event Page\n\nLive Updates Every 5-10 Seconds:\n\n✓ Poll for new matches\n✓ Auto-load results\n✓ Real-time gallery\n✓ No manual refresh"
add_box(slide, 5.3, 1.1, 4.2, 2.8, user_text, RGBColor(219, 234, 254))

# Slide 6: Tech Stack
slide = add_content_slide(prs, "Technology Stack")

tech = [
    ("Backend", "Python Flask, Gunicorn"),
    ("AI/ML", "DeepFace, Facenet512, OpenCV"),
    ("Storage", "Google Drive, MongoDB"),
    ("Frontend", "HTML5, CSS3, JavaScript"),
    ("QR Code", "qrcode[pil]"),
    ("Email", "Gmail SMTP")
]

for i, (cat, tech_name) in enumerate(tech):
    x = 0.5 if i < 3 else 5.3
    y = 1.2 + ((i % 3) * 0.95)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.2), Inches(0.75))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(243, 244, 246)
    shape.line.color.rgb = PRIMARY

    text_frame = shape.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = f"{cat}  →  {tech_name}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.LEFT

# Slide 7: Key Features
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = PRIMARY

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
title_frame = title_box.text_frame
title_frame.text = "Key Features"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

features = [
    ("⚡", "Fast", "Face match < 2 sec"),
    ("🔄", "Auto-Sync", "Real-time updates"),
    ("🌐", "Google Drive", "Direct integration"),
    ("📱", "Mobile", "Responsive design"),
    ("🔒", "Secure", "Public folders only"),
    ("📊", "Scalable", "50+ photos easily")
]

for i, (emoji, title, desc) in enumerate(features):
    row = i // 3
    col = i % 3
    x = 0.5 + (col * 3)
    y = 1.2 + (row * 1.5)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = PRIMARY

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_top = Inches(0.15)

    p = text_frame.paragraphs[0]
    p.text = f"{emoji} {title}\n{desc}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

# Slide 8: Closing
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = ACCENT

closing_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
closing_frame = closing_box.text_frame
closing_frame.text = "Ready to Find Your Photos?"
closing_frame.paragraphs[0].font.size = Pt(48)
closing_frame.paragraphs[0].font.bold = True
closing_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
closing_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
tagline_frame = tagline_box.text_frame
tagline_frame.text = "Scan • Upload • Match • Download"
tagline_frame.paragraphs[0].font.size = Pt(20)
tagline_frame.paragraphs[0].font.color.rgb = PRIMARY
tagline_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

prs.save("C:\\find_photos\\QR\\PhotoFinder_Flow_Diagram.pptx")
print("✅ Presentation created: PhotoFinder_Flow_Diagram.pptx")
